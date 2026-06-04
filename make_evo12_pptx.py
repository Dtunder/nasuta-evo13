import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_PPTX = os.path.join(HERE, "Nasuta_vs_Evo12_Presentation.pptx")

# Farben (Dark Mode Cybernetic Theme)
def rgb(hexstr): return RGBColor.from_string(hexstr.lstrip("#"))
BG = rgb("0B1117")
PANEL = rgb("131C25")
TEXT = rgb("E2E8F0")
MUTED = rgb("64748B")
ACC = rgb("38BDF8")
NAS = rgb("F43F5E")
SFULL = rgb("10B981")

def set_bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def add_text(slide, l, t, w, h, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb

def rrect(slide, l, t, w, h, fill=PANEL, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    return sp

def header(slide, kicker, title):
    add_text(slide, 0.55, 0.30, 11, 0.4, kicker, 13, ACC, bold=True)
    add_text(slide, 0.55, 0.62, 12, 0.8, title, 30, TEXT, bold=True)
    ln = slide.shapes.add_connector(2, Inches(0.6), Inches(1.45), Inches(12.7), Inches(1.45))
    ln.line.color.rgb = ACC; ln.line.width = Pt(2)

def footer(slide, n):
    add_text(slide, 0.55, 7.05, 8, 0.3, "Oekolopoly RL · Nasuta vs Evo 12", 8, MUTED)
    add_text(slide, 12.3, 7.05, 0.8, 0.3, str(n), 8, MUTED, align=PP_ALIGN.RIGHT)

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); return s

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # 1. TITLE
    s = blank(prs)
    add_text(s, 1, 2.4, 11.3, 1.0, "ÖKOLOPOLY EVO 12", 44, TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 1, 3.5, 11.3, 0.6, "Ein Kybernetischer Sieg über die Nasuta-Baseline", 22, ACC, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 1, 4.2, 11.3, 0.5, "Reinforcement Learning in hochgradig gekoppelten Systemen", 16, MUTED, align=PP_ALIGN.CENTER)
    add_text(s, 1, 6.4, 11.3, 0.4, "DeepMindMap AI & Shubham Jayswal · 2026", 11, MUTED, align=PP_ALIGN.CENTER)

    # 2. THE PROBLEM (NASUTA BASELINE)
    s = blank(prs); header(s, "DAS PROBLEM", "Die Malthusianische Falle & Nasuta Baseline")
    add_text(s, 0.55, 1.6, 12.2, 0.8, "In Ökolopoly führen positive Rückkopplungen (Bevölkerungswachstum, Produktion) unweigerlich in den Systemkollaps. Standard-RL und die Nasuta-Baseline scheitern an den späten Runden.", 13, TEXT)
    
    rrect(s, 0.6, 2.6, 5.7, 3.5)
    add_text(s, 0.85, 2.8, 5.2, 0.4, "Die Nasuta-Baseline (Paper)", 16, NAS, bold=True)
    add_text(s, 0.9, 3.4, 5.2, 2.5, "• Reine Optimierung ohne systemisches Verständnis\n• Kollabiert typischerweise in Runde 9\n• Agenten optimieren kurzfristige Gewinne\n• Folge: Überbevölkerung oder Umweltkollaps\n• Fazit: Offener Regelkreis scheitert.", 13, TEXT)

    rrect(s, 6.9, 2.6, 5.7, 3.5)
    add_text(s, 7.15, 2.8, 5.2, 0.4, "Standard RL (Evo 11)", 16, MUTED, bold=True)
    add_text(s, 7.2, 3.4, 5.2, 2.5, "• Agent nutzt Exploits ('Wall-Hugging')\n• Provoziert 'Early Suicide' (Absichtlicher Tod in Runde 6)\n• Findet lokale Optima statt globalem Überleben\n• Erkennt nicht den Wert von langfristiger Bildung", 13, TEXT)
    footer(s, 2)

    # 3. OUR SOLUTION (DANGER ZONE PENALTY)
    s = blank(prs); header(s, "UNSER ANSATZ", "Evo 12.3: Die Danger-Zone Penalty")
    add_text(s, 0.55, 1.6, 12.2, 0.8, "Um das Systemverhalten zu korrigieren, haben wir die kybernetische 'Danger-Zone Penalty' eingeführt. Sie bestraft extremales Wachstum präventiv und zwingt das Netz zur Balance.", 13, TEXT)
    
    rrect(s, 0.6, 2.6, 12.0, 1.5, fill=rgb("1E293B"))
    add_text(s, 0.85, 2.8, 11.5, 0.4, "Die quadratische Bestrafung (Python Gym Wrapper)", 14, ACC, bold=True)
    code = "population = unwrapped.V[unwrapped.POPULATION]\nif population > 30:\n    shaped_reward -= 0.1 * ((population - 30) ** 2)"
    add_text(s, 0.9, 3.2, 11.5, 0.8, code, 12, SFULL)

    add_text(s, 0.6, 4.4, 12.2, 2.0, "Wirkung:\n1. Zerschlägt das lokale Optimum des 'Early Suicides'.\n2. Zwingt den PPO-Agenten von Runde 1 an massiv in Aufklärung/Bildung zu investieren.\n3. Das exponentielle Bevölkerungswachstum wird strukturell gebrochen.", 13, TEXT)
    footer(s, 3)

    # 4. RESULTS
    s = blank(prs); header(s, "DAS ERGEBNIS", "100% Win-Rate nach 1.000.000 Steps")
    
    rrect(s, 0.6, 1.8, 12.0, 2.0, fill=SFULL)
    add_text(s, 0.85, 2.0, 11.5, 0.5, "10 / 10 Seeds überleben Runde 30", 24, BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 0.85, 2.7, 11.5, 1.0, "Durchschnittliche Runden: 30.00\nTodesursachen: KEINE (Alle 'survived')", 16, BG, align=PP_ALIGN.CENTER)

    add_text(s, 0.6, 4.2, 12.2, 1.0, "Der Phasenübergang ('Die Singularität'):", 16, ACC, bold=True)
    add_text(s, 0.6, 4.7, 12.2, 2.0, "• Nach 750k Steps: Agent stirbt in Runde 14 an verschiedenen Ursachen (Whack-a-Mole).\n• Nach 1M Steps: Das Nadelöhr (Runde 14) ist durchbrochen.\n• Sobald Runde 14 überlebt wird, nutzt der Agent die positiven Rückkopplungen (hohe Bildung = stabile Bevölkerung = hohe Lebensqualität) zu seinem Vorteil.", 14, TEXT)
    footer(s, 4)

    # 5. FORENSIC AUDIT
    s = blank(prs); header(s, "AUDIT & BEWEIS", "Echter kybernetischer Sieg (Kein Exploit)")
    add_text(s, 0.55, 1.6, 12.2, 0.8, "Forensische Analyse (Seed 42) der Endwerte in Runde 30 beweist, dass der Agent das System absolut stabilisiert hat, ohne 'Wall-Hugging' (Kleben an den Randwerten) zu betreiben.", 13, TEXT)
    
    rrect(s, 0.6, 2.5, 5.7, 3.5)
    add_text(s, 0.85, 2.7, 5.2, 0.4, "Grenzwerte & Stabilität", 16, ACC, bold=True)
    add_text(s, 0.9, 3.3, 5.2, 2.5, "• Max End-Zustand: 39 (Kritisch: >45)\n• Min End-Zustand: 10 (Kritisch: <3)\n\nAlle Sektoren oszillieren sicher im zentralen Balance-Korridor [10, 39].", 14, TEXT)

    rrect(s, 6.9, 2.5, 5.7, 3.5)
    add_text(s, 7.15, 2.7, 5.2, 0.4, "Strategie des Agenten", 16, SFULL, bold=True)
    add_text(s, 7.2, 3.3, 5.2, 2.5, "1. Früher, radikaler Aufbau von Aufklärung.\n2. Sanierung kompensiert Industrie-Schaden exakt.\n3. Politik bleibt bei 19, Bevölkerung stagniert sanft bei 30.\n\nFazit: Das Netzwerk hat die System Dynamics vollständig entschlüsselt.", 14, TEXT)
    footer(s, 5)

    prs.save(OUT_PPTX)
    print(f"Erfolgreich erstellt: {OUT_PPTX}")

if __name__ == "__main__":
    build()
