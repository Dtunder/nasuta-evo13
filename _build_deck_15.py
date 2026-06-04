# -*- coding: utf-8 -*-
"""Baut das 15-Folien-Evo-12-Deck (high-level, breites Publikum).

Quelle aller Zahlen: _EVO12_FACTS.md (Single Source of Truth).
Stil/Helper wiederverwendet aus _build_deck.py (dunkles Theme, python-pptx).
WICHTIG: quadratische Pop-Strafe = FEHLSCHLAG (Suicide-Trap Ø6), NIE als Gewinner.
Gewinner = Spread + Ceiling-Guard (Ø30 / 100 %).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "_assets")

def A(name):
    """Asset-Pfad; gibt None zurueck wenn nicht vorhanden (Bild wird dann uebersprungen)."""
    p = os.path.join(ASSET, name)
    return p if os.path.exists(p) else None

# ---------- Palette (aus _build_deck.py) ----------
BG=RGBColor(0x0E,0x17,0x26); CARD=RGBColor(0x18,0x24,0x3D); CARD2=RGBColor(0x10,0x1B,0x30)
GREEN=RGBColor(0x34,0xD3,0x99); CYAN=RGBColor(0x38,0xBD,0xF8); RED=RGBColor(0xF8,0x71,0x71)
AMBER=RGBColor(0xF5,0xB7,0x4D); VIOLET=RGBColor(0xA7,0x8B,0xFA)
TXT=RGBColor(0xE7,0xEC,0xF4); MUT=RGBColor(0x93,0xA3,0xBD)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; W=13.333

# ================= Helper (aus _build_deck.py) =================
def slide():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid()
    s.background.fill.fore_color.rgb=BG; return s

def box(s,l,t,w,h,text,size=18,color=TXT,bold=False,align=PP_ALIGN.LEFT,font='Segoe UI',
        anchor=MSO_ANCHOR.TOP,ls=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if ls: p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.name=font; r.font.color.rgb=color
    return tb

def card(s,l,t,w,h,fill=CARD,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.shadow.inherit=False
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    return sh

def bar(s,l,t,w,h,color):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background(); sh.shadow.inherit=False
    return sh

def kicker(s,text,color=CYAN):
    bar(s,0.7,0.62,0.12,0.42,color)
    box(s,0.95,0.55,11,0.5,text,15,color,bold=True)

def title(s,text):
    box(s,0.95,1.05,11.6,0.8,text,30,TXT,bold=True)

def foot(s,n):
    box(s,0.7,7.04,9,0.35,"Oekolopoly Evo 12  —  RL in gekoppelten Systemen",10,MUT)
    box(s,12.4,7.04,0.6,0.35,str(n),10,MUT,align=PP_ALIGN.RIGHT)

def pic(s,name,l,t,w):
    """Bild platzieren falls vorhanden; sonst Platzhalter-Karte."""
    p=A(name)
    if p:
        s.shapes.add_picture(p,Inches(l),Inches(t),width=Inches(w))
    else:
        card(s,l,t,w,w*0.55,CARD2,line=MUT)
        box(s,l,t+w*0.25,w,0.5,f"[{name} fehlt]",12,MUT,align=PP_ALIGN.CENTER)

def bullets(items):
    return "\n".join("•  "+b for b in items)

# ============================================================
# FOLIE 1 — TITEL
# ============================================================
s=slide(); bar(s,0,0,0.28,7.5,CYAN)
box(s,0.95,1.75,11.6,1.3,"ÖKOLOPOLY EVO 12",52,TXT,bold=True)
box(s,0.97,3.0,11.6,0.7,"Eine KI lernt, ein ganzes Land im Gleichgewicht zu halten",24,CYAN,bold=True)
box(s,0.97,3.75,11.6,0.6,"Reinforcement Learning in einem hochgradig gekoppelten Steuerungssystem",16,MUT)
chips=[("Ø 30.00","von 30 Runden",GREEN),("100 %","Win-Rate (10 Seeds)",CYAN),("1 Mio.","Trainings-Steps",AMBER)]
cx=0.95
for big,small,c in chips:
    card(s,cx,4.7,2.65,1.2,CARD,line=c)
    box(s,cx,4.85,2.65,0.6,big,28,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,cx,5.45,2.65,0.4,small,12,MUT,align=PP_ALIGN.CENTER); cx+=2.85
box(s,0.97,6.6,11.6,0.5,"Shubham Jayswal  ·  DeepMindMap AI  ·  2026",13,MUT)

# ============================================================
# FOLIE 2 — WAS IST OEKOLOPOLY?
# ============================================================
s=slide(); kicker(s,"DAS SPIEL",CYAN); title(s,"Was ist Oekolopoly?")
box(s,0.97,1.9,11.4,0.9,"Ein kybernetisches Brettspiel von Frederic Vester: Du regierst ein kleines "
    "fiktives Land und steuerst es über mehrere Runden — alle Sektoren beeinflussen sich gegenseitig.",
    15,MUT,ls=1.2)
card(s,0.95,3.05,5.7,3.3,CARD)
box(s,1.25,3.25,5.1,0.5,"8 gekoppelte Sektoren",18,CYAN,bold=True)
box(s,1.25,3.85,5.1,2.4,bullets([
    "Bevölkerung & Bevölkerungs-Wachstum","Produktion & Umwelt",
    "Lebensqualität & Sanierung","Aufklärung & Politik"]),14,TXT,ls=1.45)
card(s,6.95,3.05,5.7,3.3,CARD)
box(s,7.25,3.25,5.1,0.5,"Das Ziel",18,GREEN,bold=True)
box(s,7.25,3.85,5.1,2.4,bullets([
    "Jede Runde Aktionspunkte verteilen","Alle Sektoren im Gleichgewicht halten",
    "Weder Kollaps noch Explosion zulassen","Möglichst lange überleben — Ziel: 30 Runden"]),14,TXT,ls=1.45)
foot(s,2)

# ============================================================
# FOLIE 3 — DIE HERAUSFORDERUNG (nichtlinear, gekoppelt)
# ============================================================
s=slide(); kicker(s,"DIE HERAUSFORDERUNG",RED); title(s,"Warum ist das so schwer?")
box(s,0.97,1.9,11.4,0.9,"Das System ist nichtlinear und stark rückgekoppelt: jede Aktion wirkt "
    "verzögert und über mehrere Sektoren. Lokale Optimierung führt fast immer in den Kollaps.",
    15,MUT,ls=1.2)
pic(s,"G3_malthus.png",0.95,3.0,5.4)
card(s,6.65,3.0,6.0,3.35,CARD); bar(s,6.65,3.0,0.1,3.35,RED)
box(s,7.0,3.2,5.5,0.5,"Die Malthusianische Falle",18,RED,bold=True)
box(s,7.0,3.8,5.5,2.5,bullets([
    "Positive Rückkopplung: Bevölkerung & Produktion wachsen sich selbst hoch",
    "Wer kurzfristig optimiert, überschießt → Kollaps",
    "Spätes Endgame ist der eigentliche Killer (Runde 21–30)",
    "Zwei tödliche Wände: untere UND obere Grenze jedes Sektors"]),13.5,TXT,ls=1.35)
foot(s,3)

# ============================================================
# FOLIE 4 — BASELINES (Zahlen aus §3)
# ============================================================
s=slide(); kicker(s,"DIE BASELINES",AMBER); title(s,"Was schaffen klassische Verfahren?")
box(s,0.97,1.85,11.4,0.7,"Drei dokumentierte Referenz-Verfahren (Multi-Seed) — selbst im besten Fall "
    "(Heuristik 24/30) erreicht keine einen vollen 30-Runden-Sieg.",15,MUT,ls=1.2)
rows=[("Paper Vanilla UCT","reines MCTS","Ø 3.28","mittel (n=25)",RED),
      ("Sovereign MCTS","MCTS + Soft-Constraints","Ø 7.08","mittel (n=25)",AMBER),
      ("Heuristik","handcodiert, deterministisch","24","hoch",CYAN)]
ry=2.75
# Kopfzeile
box(s,1.0,2.55,5.0,0.4,"Methode",13,MUT,bold=True)
box(s,6.6,2.55,3.0,0.4,"Ø Runden",13,MUT,bold=True,align=PP_ALIGN.CENTER)
box(s,9.7,2.55,2.9,0.4,"Konfidenz",13,MUT,bold=True,align=PP_ALIGN.CENTER)
for name,typ,val,conf,c in rows:
    card(s,0.95,ry,11.7,0.95,CARD); bar(s,0.95,ry,0.1,0.95,c)
    box(s,1.3,ry+0.12,5.1,0.4,name,17,TXT,bold=True)
    box(s,1.3,ry+0.5,5.1,0.35,typ,12,MUT)
    box(s,6.6,ry+0.22,3.0,0.55,val,24,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,9.7,ry+0.28,2.9,0.4,conf,13,MUT,align=PP_ALIGN.CENTER)
    ry+=1.1
box(s,0.97,6.25,11.6,0.6,"Hinweis: Eine alte \"30 Runden / 100 %\"-Zahl war eine handcodierte "
    "Mogel-Heuristik, kein Lern-Erfolg. Die ehrliche Heuristik-Baseline = 24 Runden.",
    12,AMBER,ls=1.15)
foot(s,4)

# ============================================================
# FOLIE 5 — UNSER ANSATZ (RL, NICHT von null)
# ============================================================
s=slide(); kicker(s,"UNSER ANSATZ",CYAN); title(s,"RL mit MaskablePPO — nicht alles ab null")
box(s,0.97,1.9,11.4,0.95,"Wir trainieren keinen Agenten auf der grünen Wiese. Wir bauen auf der "
    "bestehenden Oekolopoly-Environment und der Vorarbeit auf — und investieren die Energie in das, "
    "was wirklich zählt: das Reward-Design.",15,MUT,ls=1.2)
card(s,0.95,3.1,5.7,3.1,CARD)
box(s,1.25,3.3,5.1,0.5,"Worauf wir aufbauen",18,CYAN,bold=True)
box(s,1.25,3.95,5.1,2.1,bullets([
    "Fertige Simulations-Environment (survival_env)",
    "Action-Masking: nur legale Züge möglich",
    "Erprobte Vorarbeit aus früheren Evo-Stufen"]),14,TXT,ls=1.45)
card(s,6.95,3.1,5.7,3.1,CARD)
box(s,7.25,3.3,5.1,0.5,"Was wir neu machen",18,GREEN,bold=True)
box(s,7.25,3.95,5.1,2.1,bullets([
    "Algorithmus: MaskablePPO",
    "Eigenes Reward-Engineering (der Kern dieser Arbeit)",
    "1 Mio. Trainings-Steps bis zur Konvergenz"]),14,TXT,ls=1.45)
foot(s,5)

# ============================================================
# FOLIE 6 — REWARD-REISE: DIE FEHLSCHLÄGE
# ============================================================
s=slide(); kicker(s,"REWARD-ENGINEERING",AMBER); title(s,"Die Reise: drei Fehlschläge zuerst")
box(s,0.97,1.85,11.4,0.7,"Bevor es funktionierte, sind drei plausible Ideen messbar gescheitert. "
    "Empirie schlägt Theorie.",15,MUT,ls=1.2)
fails=[("Quadratische Pop-Strafe","-0.1 · (pop-30)²","Ø 6","Suicide-Trap: früh sterben minimiert die Strafe",RED),
       ("Centering / Min-Margin","Abstand-zur-Mitte belohnen","Ø 6","Überpopulations-Plateau, kein Endgame-Signal",RED),
       ("Spread ohne Ceiling-Guard","nur Sektoren zusammenhalten","Ø 21","kippt an der oberen Wand (QoL 29→30)",AMBER)]
fy=2.7
for name,formel,val,why,c in fails:
    card(s,0.95,fy,11.7,1.1,CARD); bar(s,0.95,fy,0.1,1.1,c)
    box(s,1.3,fy+0.1,6.0,0.4,name,16,TXT,bold=True)
    box(s,1.3,fy+0.5,6.0,0.45,formel,12,MUT,font='Consolas')
    box(s,7.4,fy+0.3,1.7,0.55,val,22,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,9.1,fy+0.28,3.5,0.6,why,12.5,MUT,ls=1.1)
    fy+=1.25
box(s,0.97,6.5,11.6,0.5,"Lehre: Eine hart bestrafende Strafe kann ein perverses Ziel erzeugen "
    "(der Agent stirbt absichtlich). Das musste raus.",12.5,AMBER,ls=1.15)
foot(s,6)

# ============================================================
# FOLIE 7 — DER GEWINNER-REWARD (Formel §1)
# ============================================================
s=slide(); kicker(s,"DER DURCHBRUCH",GREEN); title(s,"Der Gewinner: Spread + Ceiling-Guard")
box(s,0.97,1.85,11.4,0.7,"Zwei Signale statt Bestrafung: Sektoren zur Mitte ziehen UND gezielt nur "
    "die obere Wand schützen — den einzigen Killer des Ø21-Laufs.",15,MUT,ls=1.2)
card(s,0.95,2.7,7.0,2.5,CARD2,line=GREEN)
code=("v = normalize(V[:8], Vmin, Vmax)            # 0..1\n\n"
      "reward += 1.0                               # Überleben/Step\n"
      "reward += 0.05 * balance                    # Balance-Nudge\n"
      "reward += 3.0 * (1 - (v.max()-v.min()))     # Spread → Mitte\n"
      "reward -= 8.0 * maximum(0, v-0.80).sum()    # Ceiling-Guard\n"
      "reward += 50 if Runde>=30 else -50          # Terminal")
box(s,1.25,2.92,6.5,2.1,code,12.5,GREEN,font='Consolas',ls=1.25)
card(s,8.2,2.7,4.45,2.5,CARD)
box(s,8.5,2.9,3.9,0.45,"Warum es wirkt",16,CYAN,bold=True)
box(s,8.5,3.45,3.9,1.7,bullets([
    "Guard (-8) stärker als Spread (+3): Endgame zuerst absichern",
    "Schwelle 0.80 greift nur im obersten 20 %-Band → R1–20 frei",
    "Kein Suicide-Trap: länger leben = mehr Reward"]),12.5,TXT,ls=1.3)
pic(s,"G1_reward_mechanism.png",0.95,5.35,11.7)
foot(s,7)

# ============================================================
# FOLIE 8 — PHASENÜBERGANG
# ============================================================
s=slide(); kicker(s,"DER PHASENÜBERGANG",GREEN); title(s,"Das Nadelöhr bei Runde 14")
box(s,0.97,1.8,11.4,0.7,"Lange passiert scheinbar nichts — dann kippt es. Bis ~750k Steps "
    "stagniert der Agent bei Runde 14, danach bricht er durch auf volle 30 Runden.",15,MUT,ls=1.15)
pic(s,"G5_convergence.png",1.55,2.55,10.2)
box(s,0.97,6.55,11.6,0.5,"Plateau bei R14 bis ~750k Steps · Durchbruch 14→30 zwischen 750k und 1M Steps.",
    13,MUT,align=PP_ALIGN.CENTER)
foot(s,8)

# ============================================================
# FOLIE 9 — ERGEBNISSE + VERGLEICHSTABELLE (beide Metriken)
# ============================================================
s=slide(); kicker(s,"DAS ERGEBNIS",GREEN); title(s,"Ø 30 / 100 % — beide Metriken")
stats=[("Ø 30.00","Runden (10 Seeds)",GREEN),("100 %","Win-Rate",CYAN),("0","Kollapse / 10 Seeds",AMBER)]
sy=2.05
for big,small,c in stats:
    card(s,0.95,sy,3.3,1.15,CARD,line=c)
    box(s,0.95,sy+0.1,3.3,0.6,big,28,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,0.95,sy+0.72,3.3,0.35,small,12,MUT,align=PP_ALIGN.CENTER); sy+=1.32
# Vergleichstabelle: beide Metriken
card(s,4.55,2.05,8.1,4.3,CARD)
box(s,4.85,2.2,7.6,0.45,"Vergleich gegen dokumentierte Baselines",16,CYAN,bold=True)
box(s,4.95,2.78,4.0,0.35,"Methode",12,MUT,bold=True)
box(s,9.1,2.78,1.8,0.35,"Runden",12,MUT,bold=True,align=PP_ALIGN.CENTER)
box(s,10.9,2.78,1.6,0.35,"Win-Rate (R30)",11,MUT,bold=True,align=PP_ALIGN.CENTER)
trows=[("Paper Vanilla UCT","Ø 3.28","0 %",RED),
       ("Sovereign MCTS","Ø 7.08","0 %",AMBER),
       ("Heuristik","24","0 %",CYAN),
       ("Evo 12 (MaskablePPO)","Ø 30.00","100 %",GREEN)]
ty=3.2
for name,runden,score,c in trows:
    bar(s,4.95,ty+0.05,0.08,0.55,c)
    box(s,5.15,ty,3.9,0.5,name,14,TXT,bold=(c==GREEN))
    box(s,9.1,ty,1.8,0.5,runden,16,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,10.9,ty,1.6,0.5,score,12,MUT,align=PP_ALIGN.CENTER)
    ty+=0.74
box(s,4.85,6.05,7.9,0.4,"Baselines: dokumentierte Bestwerte (n=25); keine erreicht je R30. Evo 12 auf Seeds 0–9 validiert: 30/30.",
    11,MUT)
foot(s,9)

# ============================================================
# FOLIE 10 — ROBUSTHEIT / RADAR
# ============================================================
s=slide(); kicker(s,"ROBUSTHEIT",CYAN); title(s,"Stabil über alle Seeds — kein Glückstreffer")
box(s,0.97,1.8,11.4,0.7,"Das Verhalten ist kein Einzelfall: über alle 10 Seeds bleibt der Agent "
    "innerhalb der sicheren Korridore — gleichmäßig, nicht durch Wall-Hugging.",15,MUT,ls=1.15)
pic(s,"G4_radar.png",0.95,2.7,5.85)
pic(s,"G2_robustness.png",7.05,2.7,5.6)
box(s,0.95,6.55,5.85,0.4,"Radar: alle Sektoren in sicheren Korridoren",12,MUT,align=PP_ALIGN.CENTER)
box(s,7.05,6.55,5.6,0.4,"Robustheit: konsistent über alle Seeds",12,MUT,align=PP_ALIGN.CENTER)
foot(s,10)

# ============================================================
# FOLIE 11 — TRAJEKTORIE-BEISPIEL
# ============================================================
s=slide(); kicker(s,"EIN BEISPIEL-LAUF",GREEN); title(s,"So sieht ein voller Lauf aus")
box(s,0.97,1.8,11.4,0.7,"Eine einzelne Partie (Seed 42): alle acht Sektoren bleiben in sicheren Korridoren, "
    "keiner berührt die tödliche obere oder untere Wand.",15,MUT,ls=1.15)
pic(s,"trajectory.png",1.35,2.55,10.6)
box(s,0.97,6.55,11.6,0.5,"Strategie: erst kontrolliert schrumpfen (Risiko senken), dann stabil "
    "bis Runde 30 wachsen.",13,MUT,align=PP_ALIGN.CENTER)
foot(s,11)

# ============================================================
# FOLIE 12 — EHRLICHE HYPERPARAMETER-NOTIZ
# ============================================================
s=slide(); kicker(s,"EHRLICHKEIT",AMBER); title(s,"Hyperparameter: best-found, kein bewiesenes Optimum")
card(s,0.95,2.0,11.7,1.6,CARD2,line=AMBER)
box(s,1.3,2.2,11.0,1.25,"„Wir können nicht garantieren, dass dies die global-optimalen Werte sind. "
    "Wir haben mit dieser Methode (Spread + Ceiling-Guard) die besten Ergebnisse unserer getesteten "
    "Varianten gefunden — Ø 30/30 Runden bei 100 % Win-Rate.“",16,AMBER,ls=1.25,anchor=MSO_ANCHOR.MIDDLE)
card(s,0.95,3.85,5.7,2.5,CARD)
box(s,1.25,4.05,5.1,0.45,"Was wir wissen",16,GREEN,bold=True)
box(s,1.25,4.6,5.1,1.6,bullets([
    "Reward-Struktur empirisch durch Ablation validiert",
    "Guard > Spread ist bewusst (Endgame-Schutz)",
    "Jede getestete Alternative war messbar schlechter"]),13,TXT,ls=1.3)
card(s,6.95,3.85,5.7,2.5,CARD)
box(s,7.25,4.05,5.1,0.45,"Was offen bleibt",16,AMBER,bold=True)
box(s,7.25,4.6,5.1,1.6,bullets([
    "Kein vollständiges HP-Grid trainiert (Rechenkosten)",
    "Exakte Gewichte sind best-found, nicht bewiesen",
    "Voller Beweis: 27 Kombis × ~1 Mio. Steps = Tage"]),13,TXT,ls=1.3)
foot(s,12)

# ============================================================
# FOLIE 13 — VERGLEICH VS. BASELINES (F_benchmark.png)
# ============================================================
s=slide(); kicker(s,"DER VERGLEICH",CYAN); title(s,"Evo 12 gegen alle Baselines")
box(s,0.97,1.8,11.4,0.7,"Direkt nebeneinander: das Reinforcement-Learning-Modell erreicht ein "
    "Vielfaches der besten klassischen Methode.",15,MUT,ls=1.15)
pic(s,"F_benchmark.png",1.35,2.5,10.6)
box(s,0.97,6.55,11.6,0.5,"Paper-MCTS Ø3.28 · Sovereign-MCTS Ø7.08 · Heuristik 24 · Evo 12 Ø30.00 (100 %).",
    13,MUT,align=PP_ALIGN.CENTER)
foot(s,13)

# ============================================================
# FOLIE 14 — LEARNINGS-ZUSAMMENFASSUNG
# ============================================================
s=slide(); kicker(s,"LEARNINGS",GREEN); title(s,"Was wir gelernt haben")
lessons=[("Reward-Design schlägt Modellgröße","Der Erfolg kam aus dem Belohnungssignal, nicht aus mehr Rechenpower.",CYAN),
         ("Harte Strafen sind gefährlich","Eine quadratische Strafe erzeugte einen Suicide-Trap — der Agent starb absichtlich früh.",RED),
         ("Gezielt statt pauschal","Nur die obere Wand zu schützen (Ceiling-Guard) löste den eigentlichen Endgame-Killer.",GREEN),
         ("Auf Vorarbeit aufbauen","Wir starteten nicht bei null — die bestehende Environment war das Fundament.",AMBER),
         ("Ehrlich bleiben","Best-found, kein bewiesenes Optimum — und wir benennen unsere Fehlschläge offen.",VIOLET)]
ly=2.05
for name,desc,c in lessons:
    card(s,0.95,ly,11.7,0.85,CARD); bar(s,0.95,ly,0.1,0.85,c)
    box(s,1.3,ly+0.1,4.6,0.65,name,15,c,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    box(s,6.1,ly+0.1,6.4,0.65,desc,12.5,TXT,anchor=MSO_ANCHOR.MIDDLE,ls=1.1)
    ly+=0.97
foot(s,14)

# ============================================================
# FOLIE 15 — SCHLUSS (Einzeiler)
# ============================================================
s=slide(); bar(s,0,0,0.28,7.5,GREEN)
box(s,0.95,1.7,11.6,0.9,"Das große Bild",30,CYAN,bold=True)
card(s,0.95,2.85,11.7,2.2,CARD2,line=GREEN)
box(s,1.4,3.0,10.8,1.9,"„Man musste der KI sehr viele Hilfen geben, damit sie ans Ziel kommt. "
    "Politik ist hochkomplex — und genau deshalb schwer mit KI zu automatisieren.“",
    24,TXT,bold=True,ls=1.25,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
box(s,0.97,5.4,11.6,0.6,"Ergebnis: Ø 30.00 Runden · 100 % Win-Rate · 0 Kollapse — über 10 Seeds.",
    16,GREEN,bold=True,align=PP_ALIGN.CENTER)
box(s,0.97,6.4,11.6,0.5,"Danke. Fragen?",16,MUT,align=PP_ALIGN.CENTER)

# ============================================================
OUT=os.path.join(HERE,"Evo12_Presentation_15slides.pptx")
prs.save(OUT)
print("GESPEICHERT:",OUT)
print("Folien:",len(prs.slides._sldIdLst))
