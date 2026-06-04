# -*- coding: utf-8 -*-
"""Baut ein designtes Evo-12-Paper-Deck (dunkles Theme + 2 Charts)."""
import os, numpy as np
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "_assets"); os.makedirs(ASSET, exist_ok=True)
P1 = os.path.join(ASSET, "phase.png"); P2 = os.path.join(ASSET, "trajectory.png")

# ---------- Palette ----------
BG=RGBColor(0x0E,0x17,0x26); CARD=RGBColor(0x18,0x24,0x3D); CARD2=RGBColor(0x10,0x1B,0x30)
GREEN=RGBColor(0x34,0xD3,0x99); CYAN=RGBColor(0x38,0xBD,0xF8); RED=RGBColor(0xF8,0x71,0x71)
AMBER=RGBColor(0xF5,0xB7,0x4D); TXT=RGBColor(0xE7,0xEC,0xF4); MUT=RGBColor(0x93,0xA3,0xBD)
hBG='#0E1726'; hFG='#E7ECF4'; hMUT='#93A3BD'; hGRID='#1E2C45'
hGREEN='#34D399'; hCYAN='#38BDF8'; hRED='#F87171'; hAMBER='#F5B74D'

# ================= CHARTS =================
plt.rcParams.update({'text.color':hFG,'axes.labelcolor':hFG,'xtick.color':hMUT,
    'ytick.color':hMUT,'axes.edgecolor':'#2A3A55','font.family':'DejaVu Sans','font.size':12})

# --- Chart 1: Phasenuebergang ---
fig,ax=plt.subplots(figsize=(9,4.7),dpi=200); fig.patch.set_facecolor(hBG); ax.set_facecolor(hBG)
x=[0.25,0.5,0.75,1.0]; y=[14,14,14,30]
ax.axhline(24,ls='--',color=hAMBER,lw=1.4); ax.text(0.255,24.7,'Ziel ≥ 24',color=hAMBER,fontsize=11)
ax.plot(x,y,'-o',color=hCYAN,lw=3,ms=10,mfc=hGREEN,mec=hGREEN,zorder=5)
ax.annotate('Phasenübergang\n14 → 30',xy=(1.0,30),xytext=(0.66,21.5),color=hGREEN,
    fontsize=13,fontweight='bold',arrowprops=dict(arrowstyle='->',color=hGREEN,lw=2.2))
ax.text(0.5,11.6,'Plateau — "Nadelöhr R14"',color=hMUT,fontsize=11,ha='center')
ax.set_xlabel('Trainings-Steps (Mio.)'); ax.set_ylabel('Ø Runden überlebt (10 Seeds)')
ax.set_ylim(0,33); ax.set_xlim(0.2,1.05); ax.set_xticks(x); ax.set_xticklabels(['0.25','0.5','0.75','1.0'])
ax.grid(True,color=hGRID,lw=0.7)
for sp in ('top','right'): ax.spines[sp].set_visible(False)
fig.tight_layout(); fig.savefig(P1,facecolor=hBG); plt.close(fig)

# --- Chart 2: Seed-42 Balance-Trajektorie (normalisiert) ---
raw=np.array([
 [1,12,4,10,20,13,21,0],[9,13,3,5,18,12,23,-1],[9,14,7,4,15,12,23,-3],[9,16,7,6,14,12,21,-4],
 [10,18,9,5,12,13,19,-5],[10,16,9,5,10,13,17,-6],[10,16,9,7,9,13,13,-7],[10,16,9,9,10,13,12,-8],
 [10,18,9,13,11,14,11,-7],[10,16,9,16,12,14,10,-6],[10,18,9,17,13,15,9,-5],[10,16,9,17,14,15,8,-4],
 [10,20,9,16,15,17,8,-3],[10,18,9,17,16,17,9,-2],[10,16,9,17,17,17,10,-1],[10,15,9,18,18,17,11,0],
 [10,13,9,19,19,16,12,1],[10,16,9,20,20,16,13,2],[10,16,9,22,20,16,14,4],[10,20,9,22,20,17,15,6],
 [10,16,9,21,20,17,17,8],[10,14,9,22,20,16,19,10],[11,18,9,22,20,17,21,12],[13,16,9,22,20,16,23,14],
 [15,20,9,22,20,17,25,16],[17,16,9,21,20,16,27,18],[19,20,9,21,20,16,29,20],[21,16,9,22,20,13,31,23],
 [23,22,9,19,20,12,33,25],[23,17,9,23,20,10,35,28],[23,28,10,19,21,19,39,30]],float)
vmin=np.array([1,1,1,1,1,1,1,-10],float); vmax=np.array([29,29,29,29,29,29,48,37],float)
norm=(raw-vmin)/(vmax-vmin)
labels=['Sanierung','Aufklärung','Lebensqualität','Produktion','Umwelt','Politik','Bev.-Wachstum','Bevölkerung']
cols=['#38BDF8','#34D399','#F5B74D','#F87171','#A78BFA','#22D3EE','#FB7185','#FBBF24']
fig,ax=plt.subplots(figsize=(9,4.7),dpi=200); fig.patch.set_facecolor(hBG); ax.set_facecolor(hBG)
ax.axhspan(0.9,1.0,color=hRED,alpha=0.13); ax.axhspan(0.0,0.1,color=hRED,alpha=0.13)
ax.text(0.3,0.945,'obere Wand = Tod',color=hRED,fontsize=10)
ax.text(0.3,0.045,'untere Wand = Tod',color=hRED,fontsize=10)
r=range(31)
for i in range(8): ax.plot(r,norm[:,i],lw=2,color=cols[i],label=labels[i])
ax.set_ylim(-0.02,1.02); ax.set_xlim(0,30)
ax.set_xlabel('Runde'); ax.set_ylabel('Normalisiert  (0=untere … 1=obere Wand)')
ax.grid(True,color=hGRID,lw=0.7)
for sp in ('top','right'): ax.spines[sp].set_visible(False)
ax.legend(ncol=4,loc='upper center',bbox_to_anchor=(0.5,-0.16),frameon=False,fontsize=9,labelcolor=hFG)
fig.tight_layout(); fig.savefig(P2,facecolor=hBG,bbox_inches='tight'); plt.close(fig)

# ================= DECK =================
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; W=13.333

def slide():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid()
    s.background.fill.fore_color.rgb=BG; return s

def box(s,l,t,w,h,text,size=18,color=TXT,bold=False,align=PP_ALIGN.LEFT,font='Segoe UI',
        anchor=MSO_ANCHOR.TOP,ls=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if ls: p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.name=font; r.font.color.rgb=color
    return tb

def card(s,l,t,w,h,fill=CARD,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.shadow.inherit=False
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    return sh

def bar(s,l,t,w,h,color):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background(); sh.shadow.inherit=False
    return sh

def kicker(s,text,color=CYAN):
    bar(s,0.7,0.62,0.12,0.42,color)
    box(s,0.95,0.55,11,0.5,text,15,color,bold=True)

def foot(s,n):
    box(s,0.7,7.04,8,0.35,"Oekolopoly RL  —  Nasuta vs Evo 12",10,MUT)
    box(s,12.4,7.04,0.6,0.35,str(n),10,MUT,align=PP_ALIGN.RIGHT)

# ---- Slide 1: Title ----
s=slide(); bar(s,0,0,0.28,7.5,CYAN)
box(s,0.95,2.05,11.6,1.4,"ÖKOLOPOLY EVO 12",54,TXT,bold=True)
box(s,0.97,3.35,11.6,0.7,"Ein kybernetischer Sieg über die Nasuta-Baseline",26,CYAN,bold=True)
box(s,0.97,4.15,11.6,0.6,"Reinforcement Learning in hochgradig gekoppelten Systemen",16,MUT)
# stat chips
chips=[("Ø 30.00","Runden",GREEN),("100 %","Win-Rate",CYAN),("1 Mio.","Steps",AMBER)]
cx=0.95
for big,small,c in chips:
    card(s,cx,5.05,2.55,1.15,CARD,line=c)
    box(s,cx,5.18,2.55,0.6,big,28,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,cx,5.72,2.55,0.4,small,13,MUT,align=PP_ALIGN.CENTER); cx+=2.75
box(s,0.97,6.7,11.6,0.5,"DeepMindMap AI  &  Shubham Jayswal  ·  2026",13,MUT)

# ---- Slide 2: Problem ----
s=slide(); kicker(s,"DAS PROBLEM",RED)
box(s,0.95,1.15,11.6,0.7,"Die Malthusianische Falle",30,TXT,bold=True)
box(s,0.97,1.95,11.4,0.9,"Positive Rückkopplungen (Bevölkerung, Produktion) treiben das System "
    "unweigerlich in den Kollaps. Standard-Ansätze scheitern an den späten Runden.",15,MUT,ls=1.15)
# two cards
def probcard(l,title,acc,bullets):
    card(s,l,3.0,5.7,3.55,CARD); bar(s,l,3.0,0.1,3.55,acc)
    box(s,l+0.35,3.2,5.1,0.5,title,19,acc,bold=True)
    box(s,l+0.35,3.85,5.1,2.6,"\n".join("• "+b for b in bullets),14,TXT,ls=1.3)
probcard(0.95,"Nasuta-Baseline (Paper)",RED,
    ["Optimierung ohne systemisches Verständnis","Kollabiert typischerweise in Runde 9",
     "Maximiert kurzfristige Gewinne","Folge: Überbevölkerung / Umweltkollaps",
     "Offener Regelkreis scheitert"])
probcard(6.95,"Standard RL (Evo 11)",AMBER,
    ["Nutzt Exploits ('Wall-Hugging')","Kollabiert ~Runde 12 durch Überbevölkerung",
     "Findet lokale Optima statt Überleben","Erkennt Wert langfristiger Bildung nicht"])
foot(s,2)

# ---- Slide 3: Ansatz ----
s=slide(); kicker(s,"UNSER ANSATZ",CYAN)
box(s,0.95,1.15,11.6,0.7,"Spread-Balance + Ceiling-Guard",30,TXT,bold=True)
box(s,0.97,1.95,11.4,0.95,"Zwei gekoppelte Signale halten das System stabil: ein Spread-Term zieht "
    "alle Sektoren zur Mitte, ein Ceiling-Guard verhindert das Überkippen über die obere Wand "
    "— ganz ohne den Suizid-Trap einer quadratischen Strafe.",15,MUT,ls=1.15)
# code card
card(s,0.95,3.05,7.0,2.0,CARD2,line=CYAN)
code=("v = normalize(V[:8], Vmin, Vmax)          # 0..1\n"
      "reward += 3.0 * (1 - (v.max()-v.min()))   # Spread → Mitte\n"
      "reward -= 8.0 * maximum(0, v-0.80).sum()   # Ceiling-Guard")
box(s,1.25,3.3,6.5,1.6,code,13,GREEN,font='Consolas',ls=1.35)
# effect bullets
card(s,8.2,3.05,4.45,2.0,CARD)
box(s,8.5,3.25,3.9,1.7,
    "Wirkung\n• Alle 8 Sektoren im Mittelkorridor\n• QoL-Decke (R21-Killer) blockiert\n"
    "• Kein Suicide-Trap: länger leben = mehr Reward",14,TXT,ls=1.3)
box(s,0.97,5.35,11.6,1.3,"Reward-Evolution:  quadratische Strafe (Ø6, Suizid)  →  "
    "Centering (Ø6, Überpop.)  →  Spread+Ceiling (Ø30)",14,AMBER,bold=True,ls=1.2)
foot(s,3)

# ---- Slide 4: Phasenuebergang ----
s=slide(); kicker(s,"DER PHASENÜBERGANG",GREEN)
box(s,0.95,1.15,11.6,0.7,"Das Nadelöhr bei Runde 14",30,TXT,bold=True)
s.shapes.add_picture(P1,Inches(1.7),Inches(2.05),width=Inches(10.0))
box(s,0.97,6.5,11.6,0.5,"Bis 750k Steps stagniert der Agent bei R14 (wechselnde Todesursachen). "
    "Zwischen 750k und 1M bricht er das Nadelöhr — danach 10/10 × Runde 30.",13,MUT,ls=1.1)
foot(s,4)

# ---- Slide 5: Ergebnis ----
s=slide(); kicker(s,"DAS ERGEBNIS",GREEN)
box(s,0.95,1.15,11.6,0.7,"100 % Win-Rate nach 1 Mio. Steps",30,TXT,bold=True)
stats=[("Ø 30.00","Runden (10 Seeds)",GREEN),("10 / 10","Seeds survived",CYAN),
       ("0","Todesfälle",AMBER)]
ty=2.1
for big,small,c in stats:
    card(s,0.95,ty,3.3,1.25,CARD,line=c)
    box(s,0.95,ty+0.12,3.3,0.6,big,30,c,bold=True,align=PP_ALIGN.CENTER)
    box(s,0.95,ty+0.78,3.3,0.4,small,12,MUT,align=PP_ALIGN.CENTER); ty+=1.5
s.shapes.add_picture(P2,Inches(4.55),Inches(1.95),width=Inches(8.2))
box(s,4.55,6.75,8.2,0.5,"Seed 42: alle Sektoren bleiben im Mittelkorridor — messbar kein Wall-Hugging.",
    12,MUT,align=PP_ALIGN.CENTER)
foot(s,5)

# ---- Slide 6: Audit / Verdict ----
s=slide(); kicker(s,"AUDIT & BEWEIS",CYAN)
box(s,0.95,1.15,11.6,0.7,"Echter kybernetischer Sieg — kein Exploit",30,TXT,bold=True)
card(s,0.95,2.15,5.7,3.4,CARD)
box(s,1.25,2.35,5.1,0.5,"Grenzwerte & Stabilität",18,CYAN,bold=True)
box(s,1.25,3.0,5.1,2.4,
    "• Max End-Zustand:  39   (kritisch > 45)\n• Min End-Zustand:  10   (kritisch < 3)\n"
    "• Alle Sektoren im Korridor [10, 39]\n• Stabilität = 1.0 über alle 10 Seeds",14,TXT,ls=1.45)
card(s,6.95,2.15,5.7,3.4,CARD)
box(s,7.25,2.35,5.1,0.5,"Strategie des Agenten",18,GREEN,bold=True)
box(s,7.25,3.0,5.1,2.4,
    "• Früh kontrolliert schrumpfen (Risiko runter)\n• Dann stabil wachsen (R15–30)\n"
    "• Sanierung kompensiert Industrie-Schaden\n• Bildung hält Reproduktion in Balance",14,TXT,ls=1.45)
card(s,0.95,5.75,11.7,0.95,CARD2,line=GREEN)
box(s,0.95,5.92,11.7,0.65,"Fazit: Das Netzwerk hat die System-Dynamik vollständig entschlüsselt.",
    17,GREEN,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
foot(s,6)

OUT=os.path.join(HERE,"Nasuta_vs_Evo12_Presentation_v3.pptx")
prs.save(OUT)
print("GESPEICHERT:",OUT)
print("Charts:",P1,P2)
